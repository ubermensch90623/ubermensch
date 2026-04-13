#!/usr/bin/env python3
"""Build HUG rental-deposit-guarantee workflow slide deck.

Usage:
    python3 scripts/build_pptx.py

Output: docs/assets/hug-deposit-guarantee-workflow.pptx
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Emu, Inches, Pt

OUT = Path(__file__).resolve().parents[1] / "docs" / "assets" / "hug-deposit-guarantee-workflow.pptx"

# 슬라이드 크기 16:9 (13.333 x 7.5 inch)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# 색상 팔레트
NAVY = RGBColor(0x0F, 0x2E, 0x5C)
BLUE = RGBColor(0x1F, 0x6F, 0xB5)
LIGHT = RGBColor(0xEA, 0xF2, 0xFB)
ACCENT = RGBColor(0xE8, 0x6A, 0x33)
GRAY = RGBColor(0x4A, 0x4A, 0x4A)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
RED = RGBColor(0xC6, 0x28, 0x28)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT_KO = "Noto Sans CJK KR"  # fallback chain handled by viewer; we set as primary


def set_run(run, *, size=18, bold=False, color=GRAY, font=FONT_KO):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_text_box(slide, left, top, width, height, text, *, size=18, bold=False, color=GRAY, align=None):
    from pptx.enum.text import PP_ALIGN
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line
        set_run(run, size=size, bold=bold, color=color)
        if align is not None:
            p.alignment = align
    return tb


def add_bg_band(slide, color, top, height):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, top, SLIDE_W, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp


def add_title(slide, text, *, color=NAVY):
    add_bg_band(slide, LIGHT, Inches(0), Inches(1.1))
    add_text_box(
        slide, Inches(0.5), Inches(0.28), Inches(12.3), Inches(0.7),
        text, size=28, bold=True, color=color,
    )


def add_footer(slide, page, total=10):
    add_text_box(
        slide, Inches(0.5), Inches(7.1), Inches(12.3), Inches(0.3),
        f"HUG 전세보증금반환보증 이행청구 워크플로우   |   {page}/{total}",
        size=10, color=GRAY,
    )


def new_blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank layout


def slide_cover(prs):
    s = new_blank(prs)
    add_bg_band(s, NAVY, Inches(0), SLIDE_H)
    add_text_box(
        s, Inches(0.8), Inches(2.4), Inches(11.7), Inches(1.2),
        "HUG 전세보증금반환보증\n이행청구 워크플로우",
        size=44, bold=True, color=WHITE,
    )
    add_text_box(
        s, Inches(0.8), Inches(4.1), Inches(11.7), Inches(0.6),
        "전세자금대출 세입자 — 집주인 경매 상황 대응 매뉴얼",
        size=22, color=LIGHT,
    )
    add_text_box(
        s, Inches(0.8), Inches(5.0), Inches(11.7), Inches(0.4),
        "주택도시보증공사(HUG) 공식 안내 기반 / 실무 참고용",
        size=14, color=LIGHT,
    )
    add_text_box(
        s, Inches(0.8), Inches(6.8), Inches(11.7), Inches(0.4),
        "작성일: 2026-04-13  |  본 자료는 법률 자문이 아닙니다.",
        size=11, color=LIGHT,
    )
    return s


def slide_scenario(prs):
    s = new_blank(prs)
    add_title(s, "사용자 시나리오")
    add_footer(s, 2)

    # 3-step scenario with arrows
    boxes = [
        ("전세자금대출로\n전세계약 체결", BLUE),
        ("집주인 주택이\n경매로 넘어감", ACCENT),
        ("전세보증금\n반환 불가", RED),
    ]
    x = Inches(0.8)
    y = Inches(2.0)
    w = Inches(3.5)
    h = Inches(2.5)
    gap = Inches(0.5)
    for i, (txt, color) in enumerate(boxes):
        left = Emu(int(x) + i * (int(w) + int(gap)))
        shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, y, w, h)
        shp.fill.solid()
        shp.fill.fore_color.rgb = color
        shp.line.fill.background()
        tf = shp.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = txt
        set_run(r, size=22, bold=True, color=WHITE)

    add_text_box(
        s, Inches(0.8), Inches(5.0), Inches(11.7), Inches(1.5),
        "→ HUG 전세보증금반환보증 이행청구 필요 (경매 2호 사고)\n"
        "→ 전세자금대출이 있어, 대위변제금은 대출은행 우선 상환 후 잔액만 세입자 수령",
        size=18, color=GRAY,
    )


def slide_accident_types(prs):
    s = new_blank(prs)
    add_title(s, "보증사고 유형 — 본 케이스는 2호 사고")
    add_footer(s, 3)

    # Two cards
    cards = [
        ("1호 사고", "계약 종료 후 미반환", "전세계약 해지·종료 후 1개월 내 정당한 사유 없이\n보증금을 돌려받지 못한 경우", BLUE, False),
        ("2호 사고", "경매·공매 후 미수령 (본 케이스)", "전세 기간 중 경매/공매가 실시되어,\n배당 후에도 보증금을 회수하지 못한 경우", ACCENT, True),
    ]
    for i, (t1, t2, body, color, highlight) in enumerate(cards):
        left = Emu(int(Inches(0.8)) + i * int(Inches(6.2)))
        top = Inches(1.6)
        w = Inches(5.7)
        h = Inches(4.8)
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = color
        card.line.width = Pt(3 if highlight else 1.5)
        # Title bar
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, Inches(1.0))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        bar.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        bp = bar.text_frame.paragraphs[0]
        bp.alignment = PP_ALIGN.CENTER
        br = bp.add_run()
        br.text = f"{t1}  —  {t2}"
        set_run(br, size=20, bold=True, color=WHITE)
        # Body
        add_text_box(
            s, Emu(int(left) + int(Inches(0.3))), Emu(int(top) + int(Inches(1.4))),
            Emu(int(w) - int(Inches(0.6))), Inches(3.0),
            body, size=17, color=GRAY,
        )
        if highlight:
            add_text_box(
                s, Emu(int(left) + int(Inches(0.3))), Emu(int(top) + int(Inches(3.6))),
                Emu(int(w) - int(Inches(0.6))), Inches(0.9),
                "★ 본 사용자 시나리오", size=18, bold=True, color=ACCENT,
            )


def slide_product_diff(prs):
    s = new_blank(prs)
    add_title(s, "반환보증 vs 대출상환보증 — 가입 상품 확인 필수")
    add_footer(s, 4)

    data = [
        ["구분", "반환보증", "대출상환보증 (전세금안심대출보증 등)"],
        ["보호 대상", "세입자", "대출 은행"],
        ["작동 방식", "집주인 대신 HUG가\n세입자에게 보증금 지급", "세입자 상환 불능 시\nHUG가 은행에 대신 상환"],
        ["본 시나리오와 관계", "★ 필수", "결합 상품 가입 시 병행 처리"],
    ]
    rows = len(data)
    cols = 3
    left = Inches(0.8)
    top = Inches(1.6)
    width = Inches(11.7)
    height = Inches(4.5)
    tbl = s.shapes.add_table(rows, cols, left, top, width, height).table
    tbl.columns[0].width = Inches(2.8)
    tbl.columns[1].width = Inches(4.0)
    tbl.columns[2].width = Inches(4.9)
    for r in range(rows):
        for c in range(cols):
            cell = tbl.cell(r, c)
            cell.text = ""
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = data[r][c]
            set_run(run, size=16, bold=(r == 0), color=(WHITE if r == 0 else GRAY))
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = NAVY
            elif c == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE


def slide_workflow(prs):
    s = new_blank(prs)
    add_title(s, "전체 단계별 워크플로우")
    add_footer(s, 5)

    steps = [
        "경매개시\n통지 수령",
        "HUG·은행\n사고 통지",
        "배당요구\n제출",
        "(필요 시)\n임차권등기",
        "배당 참여\n배당표 수령",
        "HUG\n이행청구",
        "HUG\n이행심사",
        "대위변제\n(은행 우선)",
        "명도\n이행금 수령",
        "HUG 구상·정산",
    ]
    # 2 rows x 5 cols
    x0 = Inches(0.5)
    y0 = Inches(1.7)
    w = Inches(2.35)
    h = Inches(1.6)
    gx = Inches(0.15)
    gy = Inches(0.9)
    for i, txt in enumerate(steps):
        row = i // 5
        col = i % 5
        if row == 1:
            col = 4 - col  # snake
        left = Emu(int(x0) + col * (int(w) + int(gx)))
        top = Emu(int(y0) + row * (int(h) + int(gy)))
        color = BLUE if i < 5 else ACCENT if i < 8 else GREEN
        shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
        shp.fill.solid()
        shp.fill.fore_color.rgb = color
        shp.line.fill.background()
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        shp.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        shp.text_frame.word_wrap = True
        p = shp.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        # number badge
        r0 = p.add_run()
        r0.text = f"{i+1}. "
        set_run(r0, size=14, bold=True, color=WHITE)
        r = p.add_run()
        r.text = txt
        set_run(r, size=15, bold=True, color=WHITE)

    add_text_box(
        s, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.7),
        "핵심 기한: 배당요구 종기 / 이행청구는 계약 종료 후 3개월 이내 / 심사·대위변제 각 약 1개월",
        size=15, bold=True, color=ACCENT,
    )


def slide_documents(prs):
    s = new_blank(prs)
    add_title(s, "HUG 제출 서류 (경매 2호 사고)")
    add_footer(s, 6)
    items = [
        "□ 보증채무이행청구서 (HUG 양식)",
        "□ 신분증 사본",
        "□ 주민등록등본 또는 초본 (말소자 포함)",
        "□ 확정일자부 전세계약서 원본",
        "□ 인감증명서 + 인감도장",
        "□ 임차권등기명령 결정문",
        "□ 임차권등기 등재된 등기부등본",
        "□ 전세계약 해지·종료 증빙 (내용증명 등)",
        "□ 배당표 사본 등 미수령액 증빙",
        "□ 채권양도계약서 + 양도통지 증빙",
        "□ (대출 있을 시) 대출 잔액증명서, 은행 상환계좌 정보, 대출약정서 사본",
    ]
    col1 = items[:6]
    col2 = items[6:]
    add_text_box(s, Inches(0.8), Inches(1.6), Inches(6.0), Inches(5.0), "\n\n".join(col1), size=17, color=GRAY)
    add_text_box(s, Inches(7.0), Inches(1.6), Inches(6.0), Inches(5.0), "\n\n".join(col2), size=17, color=GRAY)


def slide_loan(prs):
    s = new_blank(prs)
    add_title(s, "전세자금대출 연계 주의점")
    add_footer(s, 7)
    points = [
        ("① 대출은행 우선 상환", "대위변제금은 HUG → 대출은행 직접 입금.\n세입자는 대출 잔액을 제외한 차액만 수령."),
        ("② 은행·HUG 동시 통지", "이자는 대위변제일까지 계속 발생.\n경매 인지 즉시 양쪽 모두에 통지."),
        ("③ 가입 상품 확인", "반환보증 단독인지, 전세금안심대출보증\n(반환+상환 결합)인지 확인."),
        ("④ 임의 변경 금지", "경매 개시 이후 대출 중도상환·갱신 등은\n반드시 은행·HUG와 협의 후 진행."),
    ]
    for i, (title, body) in enumerate(points):
        row = i // 2
        col = i % 2
        left = Emu(int(Inches(0.8)) + col * int(Inches(6.2)))
        top = Emu(int(Inches(1.6)) + row * int(Inches(2.6)))
        w = Inches(5.7)
        h = Inches(2.3)
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT
        card.line.color.rgb = BLUE
        card.line.width = Pt(1.5)
        add_text_box(s, Emu(int(left) + int(Inches(0.3))), Emu(int(top) + int(Inches(0.2))),
                     Emu(int(w) - int(Inches(0.6))), Inches(0.6), title, size=18, bold=True, color=NAVY)
        add_text_box(s, Emu(int(left) + int(Inches(0.3))), Emu(int(top) + int(Inches(0.9))),
                     Emu(int(w) - int(Inches(0.6))), Inches(1.3), body, size=15, color=GRAY)


def slide_checklist(prs):
    s = new_blank(prs)
    add_title(s, "세입자 Do / Don't 체크리스트")
    add_footer(s, 8)

    # Do (green)
    do_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(5.9), Inches(5.2))
    do_box.fill.solid(); do_box.fill.fore_color.rgb = WHITE
    do_box.line.color.rgb = GREEN; do_box.line.width = Pt(2.5)
    add_text_box(s, Inches(1.0), Inches(1.75), Inches(5.5), Inches(0.6), "✅ 반드시 할 것", size=22, bold=True, color=GREEN)
    do_items = [
        "· 경매개시결정문 수령 즉시 HUG·은행 동시 통지",
        "· 배당요구 기한 내 제출",
        "· 전입·점유·확정일자 유지",
        "· 이사 필요 시 임차권등기 먼저",
        "· 계약서·내용증명·등기부등본 원본 보관",
        "· 계약 종료 후 3개월 이내 HUG 이행청구",
    ]
    add_text_box(s, Inches(1.0), Inches(2.5), Inches(5.5), Inches(4.2), "\n\n".join(do_items), size=15, color=GRAY)

    # Don't (red)
    dont_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(1.6), Inches(5.9), Inches(5.2))
    dont_box.fill.solid(); dont_box.fill.fore_color.rgb = WHITE
    dont_box.line.color.rgb = RED; dont_box.line.width = Pt(2.5)
    add_text_box(s, Inches(7.1), Inches(1.75), Inches(5.5), Inches(0.6), "❌ 절대 하지 말 것", size=22, bold=True, color=RED)
    dont_items = [
        "· 임차권등기 없이 전출(이사)",
        "· 경매 진행 중 점유 해제 / 주소지 변경",
        "· 배당요구 종기 경과 후 신고 시도",
        "· 임대인과 사적 합의로 권리 포기",
        "· 중요 서류 원본 분실",
        "· 청구 기한(3개월) 도과",
    ]
    add_text_box(s, Inches(7.1), Inches(2.5), Inches(5.5), Inches(4.2), "\n\n".join(dont_items), size=15, color=GRAY)


def slide_deadlines(prs):
    s = new_blank(prs)
    add_title(s, "핵심 기한 · 연락처")
    add_footer(s, 9)

    items = [
        ("⏰ 배당요구 종기", "법원 공고일 기준. 놓치면 배당 불가", ACCENT),
        ("⏰ 이행청구 기한", "전세계약 종료 후 3개월 이내", ACCENT),
        ("⏰ 이행심사·대위변제", "각 약 1개월 (합산 약 2개월)", BLUE),
        ("📞 HUG 콜센터", "1566-9009", NAVY),
        ("🌐 모바일HUG 이행청구", "onestop.khug.or.kr/webView/webBiz/excu/loahExcu001", NAVY),
        ("🌐 보증이행 안내", "khug.or.kr/hug/web/ge/er/geer001200.jsp", NAVY),
    ]
    y = Inches(1.6)
    for i, (t, v, c) in enumerate(items):
        top = Emu(int(y) + i * int(Inches(0.85)))
        # label
        shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), top, Inches(4.0), Inches(0.7))
        shp.fill.solid(); shp.fill.fore_color.rgb = c; shp.line.fill.background()
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        shp.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = shp.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = t; set_run(r, size=16, bold=True, color=WHITE)
        # value
        add_text_box(s, Inches(5.1), top, Inches(7.8), Inches(0.7), v, size=17, bold=True, color=GRAY)


def slide_refs(prs):
    s = new_blank(prs)
    add_title(s, "참고 링크 · 면책")
    add_footer(s, 10)
    refs = [
        "· HUG 전세보증금반환보증 사고란 — khug.or.kr/hug/web/ge/er/geer001100.jsp",
        "· HUG 전세보증금반환보증 이행절차 — khug.or.kr/hug/web/ge/er/geer001200.jsp",
        "· 모바일HUG 전세보증 이행청구 — onestop.khug.or.kr/webView/webBiz/excu/loahExcu001",
        "· HUG 이용절차 및 제출서류 — khug.or.kr/hug/web/ig/dr/igdr000002.jsp",
        "· HUG 전세금안심대출보증 — khug.or.kr/hug/web/ig/dl/igdl000001.jsp",
        "· 찾기쉬운 생활법령정보 — 배당요구 / 임차권등기명령 / 강제경매",
        "· 한국주택금융공사(HF) 일반전세지킴보증 — hf.go.kr/ko/sub02/sub02_05_01.do",
    ]
    add_text_box(s, Inches(0.8), Inches(1.5), Inches(12.0), Inches(4.0), "\n\n".join(refs), size=15, color=GRAY)

    add_bg_band(s, LIGHT, Inches(5.5), Inches(1.8))
    add_text_box(
        s, Inches(0.8), Inches(5.7), Inches(12.0), Inches(1.6),
        "※ 본 자료는 공개 정보 기반 실무 참고용이며 법률 자문이 아닙니다.\n"
        "   실제 대응 시 HUG 공식 안내 및 변호사·대한법률구조공단 상담을 병행하십시오.\n"
        "   서류·기한·절차는 변경될 수 있으므로 접수 전 최신 안내를 재확인 바랍니다.",
        size=14, color=NAVY,
    )


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_cover(prs)
    slide_scenario(prs)
    slide_accident_types(prs)
    slide_product_diff(prs)
    slide_workflow(prs)
    slide_documents(prs)
    slide_loan(prs)
    slide_checklist(prs)
    slide_deadlines(prs)
    slide_refs(prs)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"Saved: {OUT}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
